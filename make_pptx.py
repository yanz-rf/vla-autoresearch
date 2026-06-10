"""Convert presentation_v2.html into a native PPTX (dark theme).

Parses each <section class="slide">, extracts title / bullets / tables /
embedded base64 images & videos, and rebuilds them as PowerPoint shapes.
"""

import base64
import re
import tempfile
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x0B, 0x0E, 0x14)
FG = RGBColor(0xE8, 0xED, 0xF4)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)
VIOLET = RGBColor(0x9C, 0x6B, 0xFF)
DIM = RGBColor(0x8F, 0xA1, 0xB5)
CARD = RGBColor(0x14, 0x1A, 0x26)
TH = RGBColor(0x1D, 0x3A, 0x5F)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
tmpdir = Path(tempfile.mkdtemp(prefix="pptx_assets_"))
asset_n = 0


def save_data_uri(uri):
    global asset_n
    m = re.match(r"data:(\w+)/(\w+);base64,(.*)", uri, re.S)
    if not m:
        return None, None
    kind, ext, payload = m.group(1), m.group(2), m.group(3)
    asset_n += 1
    p = tmpdir / f"asset{asset_n}.{ext}"
    p.write_bytes(base64.b64decode(payload))
    return p, kind


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_title(slide, text, top=Inches(0.35)):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    return box


def add_bullets(slide, items, top, height, size=16):
    box = slide.shapes.add_textbox(Inches(0.7), top, SLIDE_W - Inches(1.4), height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size)
        r.font.color.rgb = FG
    return box


def html_text(el):
    t = el.get_text(" ", strip=True)
    return re.sub(r"\s+", " ", t)


def add_table(slide, table_el, top):
    rows = table_el.find_all("tr")
    ncols = max(len(r.find_all(["th", "td"])) for r in rows)
    width = SLIDE_W - Inches(1.4)
    height = Inches(0.4) * len(rows)
    shape = slide.shapes.add_table(len(rows), ncols, Inches(0.7), top, width, height)
    tbl = shape.table
    for i, row in enumerate(rows):
        cells = row.find_all(["th", "td"])
        for j in range(ncols):
            cell = tbl.cell(i, j)
            txt = html_text(cells[j]) if j < len(cells) else ""
            cell.text = txt
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].font.size = Pt(13)
                para.runs[0].font.color.rgb = FG if i else RGBColor(0xBF, 0xE3, 0xFF)
                para.runs[0].font.bold = i == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i else TH
    return top + height


def fit_picture(slide, path, top, max_h):
    from PIL import Image

    with Image.open(path) as im:
        w, h = im.size
    avail_w = SLIDE_W - Inches(1.0)
    scale = min(avail_w / Emu(w * 9525), max_h / Emu(h * 9525))
    pw, ph = Emu(int(w * 9525 * scale)), Emu(int(h * 9525 * scale))
    left = (SLIDE_W - pw) / 2
    slide.shapes.add_picture(str(path), left, top, pw, ph)
    return top + ph


html = Path("presentation_v2.html").read_text()
soup = BeautifulSoup(html, "lxml")
sections = soup.select("section.slide")
prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]

for si, sec in enumerate(sections):
    slide = prs.slides.add_slide(blank)
    add_bg(slide)

    h1 = sec.find("h1")
    h2 = sec.find("h2")
    title = html_text(h2) if h2 else (html_text(h1) if h1 else f"Slide {si + 1}")

    if h1:  # title slide: centered hero
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), SLIDE_W - Inches(1.6), Inches(4.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = html_text(h1)
        r.font.size = Pt(54)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        for sel, size, color in [(".sub", 24, FG), (".kicker", 13, DIM), (".repo", 16, VIOLET)]:
            el = sec.select_one(sel)
            if el:
                p = tf.add_paragraph()
                p.space_before = Pt(14)
                r = p.add_run()
                r.text = html_text(el)
                r.font.size = Pt(size)
                r.font.color.rgb = color
        stats = sec.select(".stat")
        if stats:
            p = tf.add_paragraph()
            p.space_before = Pt(18)
            r = p.add_run()
            r.text = "   |   ".join(
                f"{html_text(s.select_one('.n'))} {html_text(s.select_one('.l'))}" for s in stats
            )
            r.font.size = Pt(15)
            r.font.color.rgb = DIM
        continue

    add_title(slide, title)
    y = Inches(1.35)

    # stat row (teammate result slides + others)
    stats = sec.select(".stat")
    if stats:
        box = slide.shapes.add_textbox(Inches(0.7), y, SLIDE_W - Inches(1.4), Inches(0.8))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "      ".join(
            f"{html_text(s.select_one('.n'))}  {html_text(s.select_one('.l'))}" for s in stats
        )
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        y += Inches(0.8)

    # tables
    for tbl in sec.find_all("table"):
        y = add_table(slide, tbl, y) + Inches(0.25)

    # images / videos
    media = [(el.name, el.get("src", "")) for el in sec.find_all(["img", "video"])]
    n_imgs = sum(1 for k, s in media if k == "img" and s.startswith("data:"))
    for kind, src in media:
        if not src.startswith("data:"):
            continue
        path, mtype = save_data_uri(src)
        if path is None:
            continue
        if kind == "img":
            remaining = SLIDE_H - y - Inches(1.3)
            y = fit_picture(slide, path, y, remaining if n_imgs == 1 else remaining / n_imgs) + Inches(0.1)
        else:  # video
            w, h = Inches(4.6), Inches(3.45)
            lefts = [Inches(1.6), Inches(7.0)]
            left = lefts[0] if not getattr(slide, "_vid_n", 0) else lefts[1]
            slide._vid_n = getattr(slide, "_vid_n", 0) + 1
            try:
                slide.shapes.add_movie(str(path), left, y, w, h, mime_type="video/mp4")
            except Exception:
                pass

    # bullets: li, cards, callouts, paragraphs
    items = [html_text(li) for li in sec.find_all("li")]
    for card in sec.select(".card"):
        ct = card.select_one(".ct")
        head = html_text(ct) + ": " if ct else ""
        body = html_text(card).replace(html_text(ct), "", 1) if ct else html_text(card)
        items.append(head + body.strip())
    for co in sec.select(".callout, p.foot, p.center"):
        items.append(html_text(co))
    if items:
        remaining = SLIDE_H - y - Inches(0.2)
        if remaining > Inches(0.6):
            add_bullets(slide, items, y, remaining, size=14 if media or stats else 17)

prs.save("presentation_v2.pptx")
print(f"saved presentation_v2.pptx ({len(sections)} slides)")
